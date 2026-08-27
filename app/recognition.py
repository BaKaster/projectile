from __future__ import annotations

import json
import re
import stat
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import anyio
from charset_normalizer import from_bytes

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".tsv", ".json", ".jsonl", ".xml", ".yaml",
    ".yml", ".log", ".ini", ".cfg", ".sql", ".py", ".js", ".ts", ".ps1",
    ".drawio", ".excalidraw", ".svg",
}
AUDIO_EXTENSIONS = {
    ".wav", ".mp3", ".m4a", ".aac", ".ogg", ".flac", ".wma", ".webm",
    ".mp4", ".avi", ".mov", ".mkv",
}
SPREADSHEET_EXTENSIONS = {".xlsx", ".xlsm", ".xltx", ".xlt", ".xls"}
ARCHIVE_EXTENSIONS = {".zip", ".7z", ".rar"}
FORMULA_ERROR = re.compile(r"#(?:REF!|DIV/0!|VALUE!|NAME\?|N/A|NUM!|NULL!)")


class UnsupportedFormatError(RuntimeError):
    pass


class UnsafeArchiveError(RuntimeError):
    pass


@dataclass(slots=True)
class RecognitionResult:
    text: str
    tables: list[dict[str, Any]] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)


class DocumentRecognizer:
    VERSION = "hybrid-fast-office-docling-whisper/2"

    def __init__(
        self,
        *,
        whisper_model: str,
        whisper_device: str,
        whisper_compute_type: str,
        archive_max_files: int,
        archive_max_uncompressed_bytes: int,
    ) -> None:
        self.whisper_model = whisper_model
        self.whisper_device = whisper_device
        self.whisper_compute_type = whisper_compute_type
        self.archive_max_files = archive_max_files
        self.archive_max_uncompressed_bytes = archive_max_uncompressed_bytes
        self._docling_converter: Any = None
        self._whisper: Any = None

    async def recognize(self, path: Path) -> RecognitionResult:
        return await anyio.to_thread.run_sync(self._recognize_sync, path, 0)

    def _recognize_sync(self, path: Path, archive_depth: int) -> RecognitionResult:
        suffix = path.suffix.lower()
        if suffix in TEXT_EXTENSIONS:
            return self._recognize_text(path)
        if suffix == ".zip":
            return self._recognize_zip(path, archive_depth)
        if suffix in {".7z", ".rar"}:
            return self._recognize_libarchive(path, archive_depth)
        if suffix in SPREADSHEET_EXTENSIONS:
            return self._recognize_spreadsheet(path)
        if suffix == ".pdf":
            return self._recognize_pdf_text(path) or self._recognize_docling(path)
        if suffix == ".docx":
            return self._recognize_docx(path) or self._recognize_docling(path)
        if suffix == ".pptx":
            return self._recognize_pptx(path) or self._recognize_docling(path)
        if suffix in AUDIO_EXTENSIONS:
            return self._recognize_audio(path)
        if suffix == ".vsdx":
            return self._recognize_vsdx(path)
        if suffix == ".mpp":
            return self._recognize_mpp(path)
        return self._recognize_docling(path)

    @staticmethod
    def _recognize_text(path: Path) -> RecognitionResult:
        raw = path.read_bytes()
        match = from_bytes(raw).best()
        if match is None:
            raise UnsupportedFormatError(f"Не удалось определить кодировку {path.name}")
        return RecognitionResult(
            text=str(match),
            metadata={"extractor": "text", "encoding": match.encoding},
        )

    def _recognize_zip(self, path: Path, archive_depth: int) -> RecognitionResult:
        if archive_depth >= 2:
            raise UnsafeArchiveError("Превышена допустимая глубина вложенных архивов")

        chunks: list[str] = []
        tables: list[dict[str, Any]] = []
        warnings: list[str] = []
        processed = 0
        actual_total = 0
        with zipfile.ZipFile(path) as archive:
            members = [item for item in archive.infolist() if not item.is_dir()]
            total_size = sum(item.file_size for item in members)
            if len(members) > self.archive_max_files:
                raise UnsafeArchiveError("В архиве слишком много файлов")
            if total_size > self.archive_max_uncompressed_bytes:
                raise UnsafeArchiveError("Архив слишком велик после распаковки")

            with tempfile.TemporaryDirectory(prefix="projectile-archive-") as temp_dir:
                root = Path(temp_dir).resolve()
                for member in members:
                    unix_mode = member.external_attr >> 16
                    if unix_mode and stat.S_ISLNK(unix_mode):
                        raise UnsafeArchiveError("Архив содержит символическую ссылку")
                    destination = (root / member.filename).resolve()
                    if root not in destination.parents:
                        raise UnsafeArchiveError("Архив содержит небезопасный путь")
                    destination.parent.mkdir(parents=True, exist_ok=True)
                    with archive.open(member) as source, destination.open("wb") as target:
                        while block := source.read(1024 * 1024):
                            actual_total += len(block)
                            if actual_total > self.archive_max_uncompressed_bytes:
                                raise UnsafeArchiveError("Архив слишком велик после распаковки")
                            target.write(block)
                    try:
                        result = self._recognize_sync(destination, archive_depth + 1)
                    except UnsupportedFormatError as error:
                        warnings.append(str(error))
                        continue
                    processed += 1
                    chunks.append(f"\n\n## Файл из архива: {member.filename}\n\n{result.text}")
                    tables.extend(result.tables)
                    warnings.extend(result.warnings)

        return RecognitionResult(
            text="".join(chunks).strip(),
            tables=tables,
            metadata={
                "extractor": "zip",
                "archive_entries": len(members),
                "processed_entries": processed,
            },
            warnings=warnings,
        )

    def _recognize_libarchive(
        self, path: Path, archive_depth: int
    ) -> RecognitionResult:
        if archive_depth >= 2:
            raise UnsafeArchiveError("Превышена допустимая глубина вложенных архивов")
        try:
            import libarchive
        except ImportError as error:
            raise UnsupportedFormatError(
                "Для RAR/7z не установлен optional extra 'recognition'"
            ) from error

        chunks: list[str] = []
        tables: list[dict[str, Any]] = []
        warnings: list[str] = []
        processed = 0
        entry_count = 0
        unpacked_bytes = 0
        with tempfile.TemporaryDirectory(prefix="projectile-archive-") as temp_dir:
            root = Path(temp_dir).resolve()
            with libarchive.file_reader(str(path)) as archive:
                for entry in archive:
                    if stat.S_ISDIR(entry.mode):
                        continue
                    if not stat.S_ISREG(entry.mode):
                        raise UnsafeArchiveError("Архив содержит ссылку или special-файл")
                    entry_count += 1
                    unpacked_bytes += max(0, entry.size or 0)
                    if entry_count > self.archive_max_files:
                        raise UnsafeArchiveError("В архиве слишком много файлов")
                    if unpacked_bytes > self.archive_max_uncompressed_bytes:
                        raise UnsafeArchiveError("Архив слишком велик после распаковки")
                    destination = (root / entry.pathname).resolve()
                    if root not in destination.parents:
                        raise UnsafeArchiveError("Архив содержит небезопасный путь")
                    destination.parent.mkdir(parents=True, exist_ok=True)
                    written = 0
                    with destination.open("wb") as output:
                        for block in entry.get_blocks():
                            written += len(block)
                            if written > self.archive_max_uncompressed_bytes:
                                raise UnsafeArchiveError("Элемент архива слишком велик")
                            output.write(block)
                    try:
                        result = self._recognize_sync(destination, archive_depth + 1)
                    except UnsupportedFormatError as error:
                        warnings.append(str(error))
                        continue
                    processed += 1
                    chunks.append(
                        f"\n\n## Файл из архива: {entry.pathname}\n\n{result.text}"
                    )
                    tables.extend(result.tables)
                    warnings.extend(result.warnings)
        return RecognitionResult(
            text="".join(chunks).strip(),
            tables=tables,
            metadata={
                "extractor": "libarchive",
                "archive_entries": entry_count,
                "processed_entries": processed,
            },
            warnings=warnings,
        )

    @staticmethod
    def _cell_text(value: Any) -> str:
        if value is None:
            return ""
        return str(value).replace("\r", " ").replace("\n", " / ").replace("|", "\\|")

    def _recognize_spreadsheet(self, path: Path) -> RecognitionResult:
        if path.suffix.lower() in {".xls", ".xlt"}:
            return self._recognize_legacy_spreadsheet(path)
        try:
            from openpyxl import load_workbook
        except ImportError as error:
            raise UnsupportedFormatError(
                "Для Excel не установлен optional extra 'recognition'"
            ) from error

        workbook = load_workbook(
            path, read_only=True, data_only=False, keep_links=False
        )
        blocks: list[str] = []
        tables: list[dict[str, Any]] = []
        formula_errors: list[dict[str, str]] = []
        formula_count = 0
        truncated = False
        for sheet in workbook.worksheets:
            blocks.append(f"\n\n## Лист: {sheet.title}")
            rows_for_table: list[list[str]] = []
            for row_index, row in enumerate(
                sheet.iter_rows(max_row=min(sheet.max_row, 5000), max_col=min(sheet.max_column, 100)),
                start=1,
            ):
                values = [self._cell_text(cell.value) for cell in row]
                while values and not values[-1]:
                    values.pop()
                if not any(values):
                    continue
                blocks.append(" | ".join(values))
                if len(rows_for_table) < 500:
                    rows_for_table.append(values)
                for cell in row:
                    value = str(cell.value or "")
                    if cell.data_type == "f" or value.startswith("="):
                        formula_count += 1
                        if FORMULA_ERROR.search(value):
                            formula_errors.append(
                                {"sheet": sheet.title, "cell": cell.coordinate, "formula": value}
                            )
            if sheet.max_row > 5000 or sheet.max_column > 100:
                truncated = True
            tables.append({"sheet": sheet.title, "rows": rows_for_table})
        workbook.close()
        warnings = []
        if truncated:
            warnings.append("Очень большой лист был ограничен при распознавании")
        if formula_errors:
            warnings.append("В книге найдены формулы с ошибочными ссылками")
        return RecognitionResult(
            text="\n".join(blocks).strip(),
            tables=tables,
            metadata={
                "extractor": "openpyxl",
                "sheet_count": len(workbook.sheetnames),
                "formula_count": formula_count,
                "formula_errors": formula_errors[:200],
            },
            warnings=warnings,
        )

    def _recognize_legacy_spreadsheet(self, path: Path) -> RecognitionResult:
        try:
            import xlrd
        except ImportError as error:
            raise UnsupportedFormatError(
                "Для старого Excel не установлен optional extra 'recognition'"
            ) from error
        workbook = xlrd.open_workbook(path, on_demand=True)
        blocks: list[str] = []
        tables: list[dict[str, Any]] = []
        for sheet in workbook.sheets():
            blocks.append(f"\n\n## Лист: {sheet.name}")
            rows: list[list[str]] = []
            for row_index in range(min(sheet.nrows, 5000)):
                values = [self._cell_text(value) for value in sheet.row_values(row_index)[:100]]
                while values and not values[-1]:
                    values.pop()
                if not any(values):
                    continue
                blocks.append(" | ".join(values))
                if len(rows) < 500:
                    rows.append(values)
            tables.append({"sheet": sheet.name, "rows": rows})
        workbook.release_resources()
        return RecognitionResult(
            text="\n".join(blocks).strip(),
            tables=tables,
            metadata={"extractor": "xlrd", "sheet_count": workbook.nsheets},
        )

    @staticmethod
    def _recognize_vsdx(path: Path) -> RecognitionResult:
        import xml.etree.ElementTree as element_tree

        pages: list[str] = []
        with zipfile.ZipFile(path) as archive:
            names = sorted(
                name
                for name in archive.namelist()
                if name.startswith("visio/pages/page") and name.endswith(".xml")
            )
            for name in names:
                root = element_tree.fromstring(archive.read(name))
                labels: list[str] = []
                for node in root.iter():
                    if node.tag.rsplit("}", 1)[-1] == "Text":
                        value = "".join(node.itertext()).strip()
                        if value:
                            labels.append(value)
                if labels:
                    pages.append(f"## Схема {Path(name).stem}\n" + "\n".join(labels))
        if not pages:
            raise UnsupportedFormatError("В VSDX не найдено текстовых подписей")
        return RecognitionResult(
            text="\n\n".join(pages),
            metadata={"extractor": "vsdx-xml", "page_count": len(names)},
        )

    @staticmethod
    def _recognize_mpp(path: Path) -> RecognitionResult:
        try:
            import jpype
            import mpxj  # noqa: F401 - registers the MPXJ classpath
        except ImportError as error:
            raise UnsupportedFormatError(
                "Для Microsoft Project не установлен optional extra 'recognition'"
            ) from error
        if not jpype.isJVMStarted():
            jpype.startJVM()
        from org.mpxj.reader import UniversalProjectReader

        project = UniversalProjectReader().read(str(path))
        lines = ["ID | Задача | Начало | Окончание | Длительность | Выполнено"]
        task_count = 0
        for task in project.getTasks():
            name = str(task.getName() or "").strip()
            if not name:
                continue
            task_count += 1
            values = [
                task.getID(), task.getName(), task.getStart(), task.getFinish(),
                task.getDuration(), task.getPercentageComplete(),
            ]
            lines.append(" | ".join(str(value or "") for value in values))
        return RecognitionResult(
            text="\n".join(lines),
            metadata={"extractor": "mpxj", "task_count": task_count},
        )

    @staticmethod
    def _recognize_pdf_text(path: Path) -> RecognitionResult | None:
        """Use the PDF text layer when it is complete enough; scans fall back to OCR."""
        try:
            import pymupdf
        except ImportError:
            return None

        try:
            document = pymupdf.open(path)
            page_count = document.page_count
            if page_count == 0 or page_count > 1000:
                document.close()
                return None

            pages: list[str] = []
            text_page_count = 0
            character_count = 0
            for page_number, page in enumerate(document, start=1):
                page_text = page.get_text("text", sort=True).strip()
                if page_text:
                    text_page_count += 1
                    character_count += len(page_text)
                    pages.append(f"## Страница {page_number}\n{page_text}")
            document.close()

            coverage = text_page_count / page_count
            minimum_characters = max(300, page_count * 60)
            if coverage < 0.8 or character_count < minimum_characters:
                return None
            return RecognitionResult(
                text="\n\n".join(pages),
                metadata={
                    "extractor": "pymupdf-text",
                    "page_count": page_count,
                    "text_page_count": text_page_count,
                    "character_count": character_count,
                },
            )
        except Exception:  # noqa: BLE001 - malformed PDFs must fall back to OCR
            return None

    @staticmethod
    def _recognize_docx(path: Path) -> RecognitionResult | None:
        """Extract native Word text and tables without loading the ML pipeline."""
        try:
            from docx import Document as WordDocument
        except ImportError:
            return None

        try:
            document = WordDocument(path)
            blocks: list[str] = []
            tables: list[dict[str, Any]] = []

            for paragraph in document.paragraphs:
                text = paragraph.text.strip()
                if text:
                    blocks.append(text)

            for section in document.sections:
                for paragraph in section.header.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        blocks.append(f"[Верхний колонтитул] {text}")
                for paragraph in section.footer.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        blocks.append(f"[Нижний колонтитул] {text}")

            for table_index, table in enumerate(document.tables):
                rows = [
                    [cell.text.strip().replace("\n", " / ") for cell in row.cells]
                    for row in table.rows
                ]
                tables.append({"index": table_index, "rows": rows})
                blocks.append(
                    f"\n## Таблица {table_index + 1}\n"
                    + "\n".join(" | ".join(row) for row in rows)
                )

            text = "\n\n".join(blocks).strip()
            image_count = sum(
                relationship.reltype.endswith("/image")
                for relationship in document.part.rels.values()
            )
            if not text or (image_count and len(text) < 500):
                return None
            return RecognitionResult(
                text=text,
                tables=tables,
                metadata={
                    "extractor": "python-docx",
                    "paragraph_count": len(document.paragraphs),
                    "table_count": len(document.tables),
                    "image_count": image_count,
                },
            )
        except Exception:  # noqa: BLE001 - malformed DOCX must fall back to Docling
            return None

    @staticmethod
    def _recognize_pptx(path: Path) -> RecognitionResult | None:
        """Extract native PowerPoint text and tables; image-only decks use OCR."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return None

        try:
            presentation = Presentation(path)
            blocks: list[str] = []
            tables: list[dict[str, Any]] = []
            image_count = 0
            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_blocks: list[str] = []
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1
                    if getattr(shape, "has_text_frame", False):
                        text = shape.text.strip()
                        if text:
                            slide_blocks.append(text)
                    if getattr(shape, "has_table", False):
                        rows = [
                            [cell.text.strip().replace("\n", " / ") for cell in row.cells]
                            for row in shape.table.rows
                        ]
                        tables.append({"slide": slide_index, "rows": rows})
                        slide_blocks.append("\n".join(" | ".join(row) for row in rows))
                if slide_blocks:
                    blocks.append(
                        f"## Слайд {slide_index}\n" + "\n\n".join(slide_blocks)
                    )

            text = "\n\n".join(blocks).strip()
            if not text or (image_count and len(text) < 500):
                return None
            return RecognitionResult(
                text=text,
                tables=tables,
                metadata={
                    "extractor": "python-pptx",
                    "slide_count": len(presentation.slides),
                    "table_count": len(tables),
                    "image_count": image_count,
                },
            )
        except Exception:  # noqa: BLE001 - malformed PPTX must fall back to Docling
            return None

    def _get_docling_converter(self):
        if self._docling_converter is not None:
            return self._docling_converter
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import (
                PdfPipelineOptions,
                TesseractCliOcrOptions,
            )
            from docling.document_converter import (
                DocumentConverter,
                ImageFormatOption,
                PdfFormatOption,
            )
        except ImportError as error:
            raise UnsupportedFormatError(
                "Для этого формата не установлен optional extra 'recognition'"
            ) from error

        pdf_options = PdfPipelineOptions()
        pdf_options.do_ocr = True
        pdf_options.do_table_structure = True
        pdf_options.ocr_options = TesseractCliOcrOptions(lang=["rus", "eng"])
        self._docling_converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pdf_options),
                InputFormat.IMAGE: ImageFormatOption(pipeline_options=pdf_options),
            }
        )
        return self._docling_converter

    def _recognize_docling(self, path: Path) -> RecognitionResult:
        try:
            conversion = self._get_docling_converter().convert(
                path, max_file_size=128 * 1024 * 1024, max_num_pages=1000
            )
            document = conversion.document
            text = document.export_to_markdown()
            table_data: list[dict[str, Any]] = []
            for index, table in enumerate(document.tables):
                try:
                    records = table.export_to_dataframe(doc=document).fillna("").to_dict(
                        orient="records"
                    )
                    records = json.loads(json.dumps(records, default=str))
                    table_data.append({"index": index, "rows": records})
                except Exception as error:  # noqa: BLE001 - table remains in Markdown
                    table_data.append({"index": index, "error": str(error)})
            return RecognitionResult(
                text=text,
                tables=table_data,
                metadata={
                    "extractor": "docling+tesseract",
                    "table_count": len(document.tables),
                },
            )
        except UnsupportedFormatError:
            raise
        except Exception as error:
            raise UnsupportedFormatError(
                f"Не удалось распознать формат {path.suffix or '<без расширения>'}: {error}"
            ) from error

    def _get_whisper(self):
        if self._whisper is not None:
            return self._whisper
        try:
            from faster_whisper import WhisperModel
        except ImportError as error:
            raise UnsupportedFormatError(
                "Для аудио не установлен optional extra 'recognition'"
            ) from error
        self._whisper = WhisperModel(
            self.whisper_model,
            device=self.whisper_device,
            compute_type=self.whisper_compute_type,
        )
        return self._whisper

    def _recognize_audio(self, path: Path) -> RecognitionResult:
        segments, info = self._get_whisper().transcribe(
            str(path), language="ru", vad_filter=True, word_timestamps=False
        )
        segment_rows: list[dict[str, Any]] = []
        lines: list[str] = []
        for segment in segments:
            line = segment.text.strip()
            if not line:
                continue
            lines.append(f"[{segment.start:.2f}-{segment.end:.2f}] {line}")
            segment_rows.append(
                {"start": segment.start, "end": segment.end, "text": line}
            )
        return RecognitionResult(
            text="\n".join(lines),
            metadata={
                "extractor": "faster-whisper",
                "language": info.language,
                "duration_seconds": info.duration,
                "segments": segment_rows,
            },
        )


def metadata_as_table(result: RecognitionResult) -> list[dict[str, Any]]:
    """Keep recognition metadata in the existing JSONB table without a migration."""
    payload = [*result.tables, {"_recognition_metadata": result.metadata}]
    return json.loads(json.dumps(payload, default=str))
